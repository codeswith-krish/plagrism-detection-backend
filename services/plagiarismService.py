import sys
import json
import re

import os
import tempfile
import platform
import subprocess

try:
    from pptx import Presentation
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    import PyPDF2
except ImportError:
    pass

class PPTConverter:
    def __init__(self):
        self.powerpoint = None
        
    def init_com(self):
        if platform.system() == 'Windows' and self.powerpoint is None:
            import comtypes.client
            import pythoncom
            pythoncom.CoInitialize()
            self.powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            # Set to True/1 for Office apps usually to avoid background hanging, or 0 if it works.
            # Using 1 works more reliably across different Office versions with comtypes.
            self.powerpoint.Visible = 1
            
    def close_com(self):
        if self.powerpoint:
            self.powerpoint.Quit()
            self.powerpoint = None
            import pythoncom
            pythoncom.CoUninitialize()

    def convert(self, filepath, temp_dir):
        base_name = os.path.splitext(os.path.basename(filepath))[0]
        
        if platform.system() == 'Windows':
            self.init_com()
            temp_filepath = os.path.join(temp_dir, f"{base_name}_converted.pptx")
            try:
                prs = self.powerpoint.Presentations.Open(os.path.abspath(filepath), WithWindow=False)
                # 24 is format for pptx (ppSaveAsOpenXMLPresentation)
                prs.SaveAs(os.path.abspath(temp_filepath), 24)
                prs.Close()
                return temp_filepath
            except Exception as e:
                raise Exception(f"Failed to convert .ppt on Windows using comtypes: {e}")
        else:
            # Linux using LibreOffice unoconv
            temp_filepath = os.path.join(temp_dir, f"{base_name}.pptx")
            command = ['libreoffice', '--headless', '--convert-to', 'pptx', '--outdir', temp_dir, filepath]
            try:
                subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            except Exception:
                command[0] = 'soffice'
                try:
                    subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                except Exception as e:
                    raise Exception(f"Failed to convert .ppt on Linux using LibreOffice: {e}")
            
            if not os.path.exists(temp_filepath):
                raise Exception("Conversion succeeded but converted file not found.")
                
            return temp_filepath

def clean_text(text):
    text = text.lower()
    text = re.sub(r'[^a-z0-9\s.]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def extract_text_single(filepath, converter):
    is_temp = False
    current_filepath = filepath
    
    try:
        if current_filepath.lower().endswith('.pdf'):
            text_list = []
            with open(current_filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text_list.append(page.extract_text() or '')
            return clean_text(" ".join(text_list))

        if current_filepath.lower().endswith('.ppt'):
            temp_dir = tempfile.gettempdir()
            current_filepath = converter.convert(current_filepath, temp_dir)
            is_temp = True
            
        prs = Presentation(current_filepath)
        text_list = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_list.append(shape.text)
        
        full_text = " ".join(text_list)
        return clean_text(full_text)
    except Exception as e:
        raise Exception(f"EXTRACTION_FAILED: {str(e)}")
    finally:
        if is_temp and os.path.exists(current_filepath):
            try:
                os.remove(current_filepath)
            except Exception:
                pass

def extract_text(filepath):
    converter = PPTConverter()
    try:
        return extract_text_single(filepath, converter)
    finally:
        converter.close_com()

def main():
    if len(sys.argv) > 1 and sys.argv[1] == "extract":
        filepath = sys.argv[2]
        try:
            print(extract_text(filepath))
        except Exception as e:
            print(str(e), file=sys.stderr)
            sys.exit(1)
        return
        
    if len(sys.argv) > 1 and sys.argv[1] == "extract_batch":
        try:
            input_data = sys.stdin.read()
            filepaths = json.loads(input_data).get("filepaths", [])
            converter = PPTConverter()
            results = []
            for fp in filepaths:
                try:
                    text = extract_text_single(fp, converter)
                    results.append({"filepath": fp, "text": text, "error": None})
                except Exception as e:
                    results.append({"filepath": fp, "text": "", "error": str(e)})
            converter.close_com()
            print(json.dumps({"results": results}))
        except Exception as e:
            print(str(e), file=sys.stderr)
            sys.exit(1)
        return

    # Compare via stdin
    try:
        input_data = sys.stdin.read()
        if not input_data:
            return
            
        data = json.loads(input_data)
        current_text = data.get("currentText", "")
        existing_texts = data.get("existingTexts", [])
        
        if not current_text or not existing_texts:
            print(json.dumps({
                "highestSimilarity": 0.0,
                "matchedProjectId": None,
                "copiedSentences": []
            }))
            return

        documents = [current_text]
        project_ids = []
        
        for item in existing_texts:
            documents.append(item.get("text", ""))
            project_ids.append(item.get("id"))
            
        if len(documents) <= 1:
            print(json.dumps({
                "highestSimilarity": 0.0,
                "matchedProjectId": None,
                "copiedSentences": []
            }))
            return

        vectorizer = TfidfVectorizer()
        tfidf_matrix = vectorizer.fit_transform(documents)
        
        cosine_sim = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:]).flatten()
        
        max_sim_index = cosine_sim.argmax()
        max_sim_score = cosine_sim[max_sim_index]
        max_sim_percentage = round(float(max_sim_score * 100), 2)
        
        matched_project_id = project_ids[max_sim_index]
        matched_text = documents[max_sim_index + 1]
        
        current_sentences = [s.strip() for s in current_text.split('.') if len(s.strip()) > 15]
        matched_sentences = [s.strip() for s in matched_text.split('.') if len(s.strip()) > 15]
        
        matched_set = set(matched_sentences)
        copied_sentences = []
        for sentence in current_sentences:
            if sentence in matched_set:
                copied_sentences.append(sentence)
                
        copied_sentences = list(set(copied_sentences))

        result = {
            "highestSimilarity": max_sim_percentage,
            "matchedProjectId": matched_project_id,
            "copiedSentences": copied_sentences
        }
        print(json.dumps(result))

    except Exception as e:
        print(json.dumps({
            "error": str(e),
            "highestSimilarity": 0.0,
            "matchedProjectId": None,
            "copiedSentences": []
        }))

if __name__ == "__main__":
    main()
